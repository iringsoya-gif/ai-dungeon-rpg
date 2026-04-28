"""AI Dungeon RPG 발표용 PPTX 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── 색상 팔레트 ──────────────────────────────────────────────────
BG      = RGBColor(0x0a, 0x0a, 0x10)   # 배경 (아주 어두운 남색)
SURFACE = RGBColor(0x11, 0x11, 0x20)   # 카드 배경
PURPLE  = RGBColor(0x9d, 0x7f, 0xe8)   # 보라 강조
GOLD    = RGBColor(0xc9, 0xa8, 0x4c)   # 금색 강조
TEXT    = RGBColor(0xe8, 0xe4, 0xf8)   # 주 텍스트
MUTED   = RGBColor(0x8a, 0x84, 0xa8)   # 보조 텍스트
RED     = RGBColor(0xef, 0x44, 0x44)   # 위험 강조
GREEN   = RGBColor(0x10, 0xb9, 0x81)   # 성공 강조
WHITE   = RGBColor(0xff, 0xff, 0xff)

W = Inches(13.33)   # 와이드스크린 16:9 폭
H = Inches(7.5)     # 와이드스크린 16:9 높이

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # 완전 빈 레이아웃


# ── 유틸 함수 ────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color=SURFACE, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=TEXT,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline(slide, lines, l, t, w, h, size=16, color=TEXT, bold=False, spacing=1.2):
    """lines: list of (text, color, bold, size) or just str"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, col, bld, sz = item, color, bold, size
        else:
            txt, col, bld, sz = item[0], item[1], item[2], item[3] if len(item) > 3 else size
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bld
        run.font.color.rgb = col
    return txBox

def add_divider(slide, t, color=PURPLE, alpha=0.4):
    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(t), Inches(12.1), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def slide_header(slide, title, subtitle=None, icon="⚔"):
    set_bg(slide)
    # 상단 accent bar
    add_rect(slide, 0, 0, 13.33, 0.08, PURPLE)
    # 아이콘 + 타이틀
    add_text(slide, icon + "  " + title,
             0.5, 0.18, 12, 0.7, size=28, bold=True, color=TEXT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.85, 12, 0.4, size=14, color=MUTED)
    add_divider(slide, 1.3)

def bullet_card(slide, l, t, w, h, title, items, title_color=PURPLE):
    add_rect(slide, l, t, w, h, SURFACE)
    add_text(slide, title, l+0.18, t+0.15, w-0.3, 0.38,
             size=15, bold=True, color=title_color)
    item_lines = []
    for item in items:
        if isinstance(item, tuple):
            item_lines.append(item)
        else:
            item_lines.append(("  • " + item, TEXT, False, 13))
    add_multiline(slide, item_lines, l+0.18, t+0.52, w-0.3, h-0.65, size=13, color=TEXT)


# ════════════════════════════════════════════════════════════════
# 슬라이드 01 — 타이틀
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)

add_rect(sl, 0, 0, 13.33, 0.1, PURPLE)
add_rect(sl, 0, 7.4, 13.33, 0.1, PURPLE)

# 배경 장식 사각형
add_rect(sl, 3.5, 1.2, 6.3, 5.1, RGBColor(0x14, 0x14, 0x24))

add_text(sl, "⚔", 6.0, 1.5, 1.5, 1.2, size=52, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(sl, "AI Dungeon RPG",
         1.5, 2.8, 10.3, 1.2, size=42, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_text(sl, "AI 게임마스터가 이끄는 인터랙티브 웹소설 RPG 서비스",
         2.0, 3.9, 9.3, 0.7, size=18, color=MUTED, align=PP_ALIGN.CENTER)
add_divider(sl, 4.75, GOLD)
add_text(sl, "202212577  박성안  |  부산대학교  |  2026.04",
         3.0, 5.0, 7.3, 0.5, size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(sl, "FastAPI  ·  React  ·  Groq LLaMA 3.3 70B  ·  SSE  ·  Vercel / Railway",
         2.5, 5.5, 8.3, 0.45, size=12, color=RGBColor(0x5a,0x54,0x78), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 슬라이드 02 — 목차
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "목차", "10분 발표 구성")

items = [
    ("01", "프로젝트 개요 및 동기",    PURPLE),
    ("02", "시스템 아키텍처 & 기술 스택", GOLD),
    ("03", "핵심 기술: AI GM · SSE · 컨텍스트 관리", PURPLE),
    ("04", "주요 기능 소개",            GOLD),
    ("05", "수업 내용 활용",            PURPLE),
    ("06", "기존 AI 챗봇과의 비교",     GOLD),
    ("07", "기술적 도전과 해결",        PURPLE),
    ("08", "결론 & Q&A",               GOLD),
]
for i, (num, label, col) in enumerate(items):
    row = i // 2
    col_x = 0.6 if i % 2 == 0 else 7.0
    add_rect(sl, col_x, 1.5 + row*1.3, 5.9, 1.1, SURFACE)
    add_text(sl, num, col_x+0.15, 1.55 + row*1.3, 0.7, 0.5,
             size=22, bold=True, color=col)
    add_text(sl, label, col_x+0.75, 1.6 + row*1.3, 5.0, 0.6,
             size=15, bold=True, color=TEXT)


# ════════════════════════════════════════════════════════════════
# 슬라이드 03 — 프로젝트 개요 & 동기
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "프로젝트 개요 & 동기", "왜 만들었는가?")

add_text(sl, "\"AI와 사람이 함께 만드는 이야기\"",
         0.5, 1.45, 12.3, 0.6, size=20, bold=True, color=GOLD, italic=True)

bullet_card(sl, 0.5, 2.1, 5.8, 2.5, "🎯  무엇을 만들었나?", [
    "AI 게임마스터(GM)가 실시간으로 서사를 생성하는 RPG",
    "플레이어의 행동·대사에 반응하는 인터랙티브 스토리텔링",
    "완전 배포된 웹 서비스 (실제 접속 가능)",
])
bullet_card(sl, 6.5, 2.1, 6.3, 2.5, "💡  만든 이유", [
    "ChatGPT로 RPG를 해보면 매번 상태를 다시 설명해야 함",
    "캐릭터 HP, 인벤토리가 자동 추적되지 않음",
    "\"AI가 진짜 GM 역할을 할 수 있을까?\" 라는 호기심",
    "수업에서 배운 API/웹개발 기술을 종합 활용하는 프로젝트",
])

add_text(sl, "🌐  실제 서비스: https://ai-dungeon-rpg.vercel.app",
         0.5, 4.75, 12.3, 0.5, size=14, color=PURPLE, bold=True)


# ════════════════════════════════════════════════════════════════
# 슬라이드 04 — 시스템 아키텍처
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "시스템 아키텍처", "전체 기술 스택 한눈에")

# 3열 구조: Frontend | Backend | AI/DB
cards = [
    ("Frontend\n(Vercel)", PURPLE, [
        "React 18 + Vite",
        "Tailwind CSS",
        "Zustand (상태관리)",
        "SSE 실시간 스트리밍",
        "Web Audio API BGM",
        "React Router v6",
    ]),
    ("Backend\n(Railway)", GOLD, [
        "FastAPI (Python)",
        "SQLAlchemy ORM",
        "SQLite DB",
        "JWT + Google OAuth",
        "Polar.sh 결제 연동",
        "Docker 컨테이너",
    ]),
    ("AI / 외부 서비스", GREEN, [
        "Groq API (무료 티어)",
        "LLaMA 3.3 70B — GM 서사",
        "LLaMA 3.3 70B — 오프닝",
        "Google OAuth 2.0",
        "Polar.sh (결제·플랜 동기화)",
        "Vercel + Railway 자동 배포",
    ]),
]
for i, (title, col, items) in enumerate(cards):
    x = 0.5 + i * 4.3
    add_rect(sl, x, 1.5, 3.9, 5.6, SURFACE)
    add_rect(sl, x, 1.5, 3.9, 0.5, col)
    add_text(sl, title, x+0.1, 1.52, 3.7, 0.46,
             size=14, bold=True, color=RGBColor(0x0a,0x0a,0x10), align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(sl, "▸  " + item,
                 x+0.18, 2.15+j*0.82, 3.55, 0.7, size=13, color=TEXT)

# 화살표 텍스트
add_text(sl, "→", 4.45, 4.1, 0.4, 0.5, size=24, color=MUTED, align=PP_ALIGN.CENTER)
add_text(sl, "→", 8.75, 4.1, 0.4, 0.5, size=24, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 슬라이드 05 — 핵심 기술 ① AI 게임마스터
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "핵심 기술 ①: AI 게임마스터", "LLM을 GM으로 활용하는 방법")

bullet_card(sl, 0.5, 1.5, 6.1, 5.6, "🤖  AI GM 동작 원리", [
    "Groq API로 LLaMA 3.3 70B 모델 호출",
    "시스템 프롬프트 = GM의 '인격'과 '규칙서'",
    "  • 세계관 설명 동적 삽입",
    "  • 현재 캐릭터 JSON 전달",
    "  • 한국 웹소설 문체 지시",
    "  • 장면 연속성 & 언어 규칙",
    "응답 마지막에 JSON 블록 포함",
    "  → 상태 변화를 구조적으로 추출",
])

# JSON 예시 코드 블록
add_rect(sl, 6.8, 1.5, 6.0, 5.6, RGBColor(0x0d, 0x0d, 0x1a))
add_text(sl, "GM 응답 형식 (구조화된 출력)",
         6.95, 1.6, 5.8, 0.4, size=12, color=PURPLE, bold=True)
code = """{
  "state_changes": {
    "hp_change": -15,
    "xp_gain": 30,
    "inventory_add": ["낡은 단검"],
    "quest_add": ["어둠의 마법사 처치"],
    "location": "어두운 지하 감옥",
    "in_battle": true
  },
  "game_over": false
}"""
add_text(sl, code, 6.95, 2.05, 5.7, 4.8,
         size=12, color=GREEN, bold=False)


# ════════════════════════════════════════════════════════════════
# 슬라이드 06 — 핵심 기술 ② SSE 실시간 스트리밍
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "핵심 기술 ②: 실시간 스트리밍 (SSE)", "ChatGPT처럼 타이핑되며 나타나는 서술")

# 흐름도
steps = [
    ("① 플레이어\n행동 입력", PURPLE),
    ("② FastAPI\nSSE 응답 시작", GOLD),
    ("③ Groq API\n스트리밍 청크", GREEN),
    ("④ 화면에\n실시간 표시", PURPLE),
    ("⑤ 완료 시\n상태 자동 반영", GOLD),
]
for i, (label, col) in enumerate(steps):
    x = 0.4 + i * 2.55
    add_rect(sl, x, 1.55, 2.2, 1.3, col)
    add_text(sl, label, x+0.08, 1.65, 2.05, 1.1,
             size=13, bold=True, color=RGBColor(0x0a,0x0a,0x10), align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(sl, "→", x+2.22, 2.0, 0.3, 0.5, size=18, color=MUTED)

bullet_card(sl, 0.5, 3.1, 5.8, 3.7, "📡  SSE 구현 포인트", [
    "FastAPI StreamingResponse + async generator",
    "이벤트: text | done | error 세 종류",
    "'done' 이벤트에 상태 변화 JSON 포함",
    "프론트: EventSource 대신 fetch + ReadableStream",
    "스트리밍 중 JSON 블록은 화면에서 숨김",
])
bullet_card(sl, 6.5, 3.1, 6.3, 3.7, "⚙️  컨텍스트 관리", [
    "8,000 토큰 초과 시 자동 압축 (요약 생성)",
    "최근 10턴만 유지 + 이전 요약 prepend",
    "오프닝 장면은 [게임 시작] 마커로 보존",
    "캐릭터 JSON은 매 턴 시스템 프롬프트에 갱신",
    "→ 장기 플레이도 맥락 유지 가능",
])


# ════════════════════════════════════════════════════════════════
# 슬라이드 07 — 주요 기능 ①
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "주요 기능 ①: 세계관 & 캐릭터 시스템", "나만의 RPG 세계 구축")

bullet_card(sl, 0.5, 1.5, 3.9, 5.6, "🌍  세계관 템플릿", [
    "판타지 — 검과 마법의 세계",
    "SF — 우주·사이버펑크",
    "공포 — 공포·오컬트",
    "현대 — 현대 도시",
    "커스텀 — 자유 입력",
    "",
    "→ 선택한 장르에 따라",
    "  직업 목록이 자동 변경",
])
bullet_card(sl, 4.6, 1.5, 4.1, 5.6, "👤  캐릭터 시스템", [
    "20개 프리셋 직업",
    "(판타지·SF·공포·현대 각 5종)",
    "커스텀 직업 직접 입력 가능",
    "",
    "스탯: HP · MP · STR",
    "        INT · AGI · CHA",
    "",
    "레벨업 시 전체 스탯 자동 증가",
    "XP 요구량 ×1.5 증가 (누적)",
])
bullet_card(sl, 8.9, 1.5, 3.9, 5.6, "⚔  게임 시스템", [
    "인벤토리 자동 추가·제거",
    "퀘스트 추가·완료 추적",
    "상태 이상 효과 시스템",
    "전투 / 탐험 상태 전환",
    "하드코어 모드 (영구 사망)",
    "일반 모드: 사망 시 패널티",
    "  → HP ½ 회복, 아이템 1개 손실",
    "게임 종료 통계 화면",
])


# ════════════════════════════════════════════════════════════════
# 슬라이드 08 — 주요 기능 ②
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "주요 기능 ②: UI/UX & 부가 기능", "사용자 경험 설계")

bullet_card(sl, 0.5, 1.5, 5.9, 3.5, "🎨  UI/UX 디자인", [
    "다크 판타지 테마 (caveduck.io 참조)",
    "**행동** 형식으로 행동·대사 구분",
    "GM 서술에서 JSON 자동 숨김",
    "실시간 스트리밍 커서 애니메이션",
    "사이드패널: HP/MP/XP 바, 인벤토리, 퀘스트",
    "캐릭터 시트 팝업 모달",
    "모바일 대응 반응형 레이아웃",
])
bullet_card(sl, 6.6, 1.5, 6.2, 3.5, "🎵  Web Audio API BGM", [
    "외부 음원 없이 브라우저 내에서 생성",
    "탐험 모드: A단조 앰비언트 드론",
    "전투 모드: E단조, 더 높고 긴박하게",
    "게임오버: D단조, 낮고 무겁게",
    "모드 전환 시 1.5초 페이드 크로스페이드",
    "♪/♩ 버튼으로 ON/OFF, localStorage 저장",
])
bullet_card(sl, 0.5, 5.2, 5.9, 2.1, "🔐  인증 & 결제", [
    "Google OAuth 2.0 소셜 로그인",
    "JWT 토큰 기반 세션 관리",
    "Polar.sh 결제 (무료/영웅 플랜)",
])
bullet_card(sl, 6.6, 5.2, 6.2, 2.1, "☁️  배포 & 운영", [
    "Frontend: Vercel (GitHub push → 자동 배포)",
    "Backend: Railway (Docker 컨테이너)",
    "GitHub main 브랜치 푸시 → 즉시 반영",
])


# ════════════════════════════════════════════════════════════════
# 슬라이드 09 — 수업 내용 활용
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "수업 내용 활용", "AI·ML 과목에서 배운 것들을 어떻게 썼나?")

topics = [
    ("LLM & 프롬프트 엔지니어링", PURPLE, [
        "→ LLM의 토큰 기반 생성 원리 이해",
        "→ 시스템 프롬프트로 모델 행동 제어",
        "→ 출력 형식 구조화 (JSON 강제)",
        "→ Few-shot: 좋은 예/나쁜 예 제시로",
        "   문체 학습 유도",
    ]),
    ("API & 비동기 처리", GOLD, [
        "→ REST API 설계와 HTTP 메서드",
        "→ async/await & SSE 스트리밍",
        "→ FastAPI를 Week5에서 직접 실습",
        "→ API 키 관리 및 환경변수 운영",
        "→ Groq 무료 API로 비용 없이 구현",
    ]),
    ("데이터 구조 & 상태 관리", GREEN, [
        "→ 캐릭터 상태를 JSON으로 직렬화",
        "→ SQLite: 관계형 DB 스키마 설계",
        "→ 컨텍스트 윈도우 = 일종의 슬라이딩 윈도우",
        "→ 토큰 수 추정 (len/4 근사)",
        "→ 자동 요약으로 정보 압축",
    ]),
    ("웹소설 스타일 프롬프트", RED, [
        "→ NLP: 텍스트 생성 방향 제어",
        "→ 문체 transfer를 프롬프트로 구현",
        "→ Bad example / Good example 기법",
        "→ 언어 규칙 명시 (한국어 전용)",
        "→ 장면 연속성 제약 조건 설계",
    ]),
]
for i, (title, col, items) in enumerate(topics):
    row = i // 2
    cx = 0.5 if i % 2 == 0 else 6.9
    add_rect(sl, cx, 1.55 + row*2.9, 6.0, 2.7, SURFACE)
    add_rect(sl, cx, 1.55 + row*2.9, 6.0, 0.42, col)
    add_text(sl, title, cx+0.15, 1.62 + row*2.9, 5.7, 0.35,
             size=13, bold=True, color=RGBColor(0x0a,0x0a,0x10))
    item_lines = [("  " + it, TEXT, False, 12) for it in items]
    add_multiline(sl, item_lines, cx+0.15, 2.05 + row*2.9, 5.7, 2.1)


# ════════════════════════════════════════════════════════════════
# 슬라이드 10 — 기존 AI 챗봇과 비교
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "기존 AI 챗봇과의 비교", "ChatGPT로 RPG를 하면 생기는 문제들")

# 헤더 행
add_rect(sl, 0.5, 1.5, 4.8, 0.55, RGBColor(0x1e,0x1e,0x30))
add_rect(sl, 5.4, 1.5, 3.6, 0.55, RGBColor(0x3a,0x2a,0x10))
add_rect(sl, 9.1, 1.5, 3.7, 0.55, PURPLE)
add_text(sl, "비교 항목", 0.6, 1.55, 4.6, 0.45, size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(sl, "ChatGPT / Claude",   5.5, 1.55, 3.4, 0.45, size=13, bold=True, color=GOLD,   align=PP_ALIGN.CENTER)
add_text(sl, "AI Dungeon RPG",     9.2, 1.55, 3.5, 0.45, size=13, bold=True, color=RGBColor(0x0a,0x0a,0x10), align=PP_ALIGN.CENTER)

rows = [
    ("캐릭터 상태 추적",     "매번 직접 설명해야 함",       "자동 저장 & 실시간 갱신",     False),
    ("HP/MP/XP 관리",       "직접 계산해야 함",            "자동 계산 & 레벨업 처리",     False),
    ("장기 플레이 맥락",     "컨텍스트 초과 시 망각",       "자동 요약 압축으로 유지",     False),
    ("게임 진행 저장",       "브라우저 새로고침 시 초기화", "DB 영구 저장, 언제든 재접속", False),
    ("하드코어 모드",        "직접 규칙 설정 필요",         "시스템 수준 보장 (서버 처리)", False),
    ("실시간 서사 스트리밍", "기본 지원 (UI 있음)",         "커스텀 SSE 구현",             True),
    ("한국 웹소설 문체",     "일반적 문어체",               "전용 프롬프트 + Few-shot",    False),
    ("장면 반응 BGM",        "없음",                        "Web Audio API 앰비언트 생성", False),
]

for i, (item, chatgpt, ours, neutral) in enumerate(rows):
    y = 2.15 + i * 0.62
    bg = RGBColor(0x11,0x11,0x20) if i % 2 == 0 else RGBColor(0x0e,0x0e,0x1a)
    add_rect(sl, 0.5, y, 4.8, 0.58, bg)
    add_rect(sl, 5.4, y, 3.6, 0.58, bg)
    add_rect(sl, 9.1, y, 3.7, 0.58, bg)
    add_text(sl, item,   0.65, y+0.08, 4.6, 0.45, size=12, color=TEXT)
    c_col = MUTED if neutral else RGBColor(0xef,0x88,0x44)
    add_text(sl, "△  " + chatgpt, 5.5, y+0.08, 3.4, 0.45, size=11, color=c_col)
    add_text(sl, "✓  " + ours,   9.2, y+0.08, 3.5, 0.45, size=11, color=GREEN)


# ════════════════════════════════════════════════════════════════
# 슬라이드 11 — 기술적 도전과 해결
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "기술적 도전과 해결", "만들면서 만난 문제들")

problems = [
    ("🔇  GM이 응답 없이 침묵", PURPLE,
     "동기 Anthropic 클라이언트가\n비동기 이벤트 루프 블로킹",
     "AsyncGroq + 메시지 배열\nuser 첫 항목 조건 강제 적용"),
    ("🌐  API 비용 문제", GOLD,
     "Anthropic API 크레딧 소진\n무료로 계속 개발 불가",
     "Groq API (무료 티어) 전환\nLLaMA 3.3 70B 무료 사용"),
    ("🇯🇵  일본어 한자 혼입", RED,
     "모델이 한국어 서술 중\n漢字·日本語를 섞어 출력",
     "프롬프트 최상단에 언어 규칙 +\n잘못된 예시 직접 명시"),
    ("📍  장면 연속성 단절", GREEN,
     "오프닝 장면이 context_manager에서\n삭제되어 모델에 전달 안 됨",
     "assistant 메시지 삭제 →\n[게임 시작] 마커 삽입으로 대체"),
    ("📝  오프닝 서술 품질 저하", PURPLE,
     "8B 모델: 시점 혼용(나/크루/당신),\n직업 무시, 첫 장면부터 사건 폭발",
     "70B 모델 전환 + 프롬프트에\n3단계 구성·2인칭 규칙·직업 힌트 명시"),
    ("💳  결제 후 플랜 미반영", GOLD,
     "Webhook 미도달 시 DB 업데이트\n안 됨 → 플랜 여전히 free",
     "성공 페이지에서 checkout_id 검증\n+ 대시보드 플랜 동기화 버튼"),
]

# 6개 카드: 3열 × 2행
for i, (title, col, prob, sol) in enumerate(problems):
    row = i // 3
    cx = 0.4 + (i % 3) * 4.3
    add_rect(sl, cx, 1.55 + row*2.9, 3.9, 2.7, SURFACE)
    add_rect(sl, cx, 1.55 + row*2.9, 3.9, 0.38, col)
    add_text(sl, title, cx+0.12, 1.6 + row*2.9, 3.7, 0.33,
             size=11, bold=True, color=RGBColor(0x0a,0x0a,0x10))
    add_text(sl, "문제: " + prob,
             cx+0.12, 2.02 + row*2.9, 3.7, 0.88,
             size=10, color=RGBColor(0xef,0x88,0x44))
    add_rect(sl, cx+0.1, 2.93 + row*2.9, 3.7, 0.02, MUTED)
    add_text(sl, "해결: " + sol,
             cx+0.12, 2.98 + row*2.9, 3.7, 0.92,
             size=10, color=GREEN)


# ════════════════════════════════════════════════════════════════
# 슬라이드 12 — 결론 & Q&A
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.1, PURPLE)
add_rect(sl, 0, 7.4, 13.33, 0.1, PURPLE)

add_text(sl, "⚔", 6.0, 0.8, 1.3, 1.0,
         size=40, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(sl, "결론",
         2.5, 1.8, 8.3, 0.9, size=36, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
add_divider(sl, 2.75, GOLD)

summary = [
    ("✦", "LLM API + 웹 기술을 조합해 실제 서비스를 처음부터 끝까지 구현"),
    ("✦", "프롬프트 엔지니어링이 AI 출력 품질을 결정적으로 좌우함을 직접 경험"),
    ("✦", "ChatGPT와의 차별점: 지속 상태·자동 추적·장면 반응 BGM"),
    ("✦", "기술적 문제(언어 혼입, 장면 단절)를 반복 실험으로 해결"),
    ("✦", "배포까지 완성 → 실제 접속 가능한 서비스"),
]
for i, (icon, text) in enumerate(summary):
    add_rect(sl, 1.5, 3.0 + i*0.78, 10.3, 0.65, SURFACE)
    add_text(sl, icon, 1.65, 3.07 + i*0.78, 0.5, 0.55, size=14, color=PURPLE, bold=True)
    add_text(sl, text, 2.25, 3.07 + i*0.78, 9.4, 0.55, size=13, color=TEXT)

add_text(sl, "🌐  https://ai-dungeon-rpg.vercel.app",
         3.0, 7.0, 7.3, 0.4, size=13, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

add_text(sl, "Q & A",
         5.5, 6.4, 2.3, 0.5, size=22, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════
# 저장
# ════════════════════════════════════════════════════════════════
out = r"C:\Users\USER\AIandMLcourse\ai-dungeon-rpg\AI_Dungeon_RPG_발표.pptx"
prs.save(out)
print(f"✓ 저장 완료: {out}")
