"""
AI Dungeon RPG — 발표용 PPT 생성기 (Phase 2 완성판)
실행: uv run python presentation/generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── 색상 팔레트 ──────────────────────────────────────────────
BG      = RGBColor(0x0a, 0x0a, 0x10)
SURFACE = RGBColor(0x11, 0x11, 0x20)
PURPLE  = RGBColor(0x9d, 0x7f, 0xe8)
PURPLE2 = RGBColor(0x6a, 0x3f, 0xa0)
GOLD    = RGBColor(0xc9, 0xa8, 0x4c)
GOLD2   = RGBColor(0x8a, 0x68, 0x20)
TEXT    = RGBColor(0xe8, 0xe4, 0xf8)
MUTED   = RGBColor(0x6b, 0x6b, 0x90)
RED     = RGBColor(0xef, 0x44, 0x44)
GREEN   = RGBColor(0x10, 0xb9, 0x81)
BORDER  = RGBColor(0x1e, 0x1e, 0x30)
BLUE    = RGBColor(0x60, 0xa5, 0xfa)

W = Inches(16)
H = Inches(9)
TOTAL = 13

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── 헬퍼 ────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def box(slide, l, t, w, h, color=SURFACE, line=BORDER, line_w=1):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h,
        size=20, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False,
        font="Malgun Gothic"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def multi_line(slide, lines, l, t, w, h,
               size=20, bold=False, color=TEXT,
               align=PP_ALIGN.LEFT, italic=False,
               font="Malgun Gothic", spacing=1.4):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb


def slide_num(slide, n):
    txt(slide, f"{n} / {TOTAL}",
        W - Inches(1.2), H - Inches(0.45),
        Inches(1.0), Inches(0.35),
        size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def accent_bar(slide, color=PURPLE):
    shape = slide.shapes.add_shape(1, 0, 0, W, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def section_tag(slide, label, color=PURPLE):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(5), Inches(0.42))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label.upper()
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Malgun Gothic"


def heading(slide, text, color=TEXT):
    txt(slide, text, Inches(0.6), Inches(0.75), Inches(14.4), Inches(0.85),
        size=34, bold=True, color=color)


def divider(slide, y, color=BORDER):
    shape = slide.shapes.add_shape(1, Inches(0.6), y, Inches(14.8), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ════════════════════════════════════════════════════════════
#  SLIDE 01 — 타이틀
# ════════════════════════════════════════════════════════════
s1 = new_slide()
accent_bar(s1, PURPLE)

# 배경 글로우
glow = s1.shapes.add_shape(9, Inches(4), Inches(1.5), Inches(8), Inches(6))
glow.fill.solid()
glow.fill.fore_color.rgb = RGBColor(0x12, 0x08, 0x26)
glow.line.fill.background()

txt(s1, "⚔  AI DUNGEON RPG",
    0, Inches(1.8), W, Inches(0.65),
    size=13, color=PURPLE, align=PP_ALIGN.CENTER, bold=True)

multi_line(s1, ["AI가 게임마스터가 되는", "웹소설 RPG"],
    0, Inches(2.55), W, Inches(2.9),
    size=62, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

txt(s1, "FastAPI  ·  React  ·  Groq Llama 3.3 70B  ·  Google OAuth  ·  Polar.sh",
    0, Inches(5.6), W, Inches(0.55),
    size=14, color=MUTED, align=PP_ALIGN.CENTER)

txt(s1, "당신의 행동에 AI가 실시간으로 반응합니다",
    0, Inches(6.3), W, Inches(0.6),
    size=17, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

txt(s1, "202212577 박성안  |  부산대학교 물리학과",
    0, H - Inches(0.65), W, Inches(0.5),
    size=12, color=MUTED, align=PP_ALIGN.CENTER)

slide_num(s1, 1)


# ════════════════════════════════════════════════════════════
#  SLIDE 02 — 문제: ChatGPT RPG의 한계
# ════════════════════════════════════════════════════════════
s2 = new_slide()
accent_bar(s2, RED)
section_tag(s2, "배경 · 동기", RED)
heading(s2, "ChatGPT로 RPG를 한다면?", RED)

problems = [
    ("❌  매번 상태를 직접 설명",
     '"HP 80, 단검 보유, 던전 3층" — 새로고침하면 사라짐'),
    ("❌  게임 시스템 없음",
     "레벨업도, 인벤토리도, 퀘스트 추적도 없는 단순 채팅"),
    ("❌  한국어 품질 문제",
     "서술 중 갑자기 한자·영어 삽입, 어색한 번역체"),
]
solutions = [
    ("✅  서버가 자동으로 상태 추적",
     "HP·XP·인벤토리·퀘스트 — DB 영구 저장"),
    ("✅  완전한 RPG 시스템",
     "레벨업·전투·NPC·하드코어 모드 내장"),
    ("✅  장르별 한국어 문체",
     "웹소설 스타일 프롬프트 + Few-shot 엔지니어링"),
]

for i, (title, sub) in enumerate(problems):
    y = Inches(2.1 + i * 1.55)
    box(s2, Inches(0.5), y, Inches(7.0), Inches(1.3),
        color=RGBColor(0x16, 0x07, 0x07), line=RGBColor(0x60, 0x18, 0x18))
    txt(s2, title, Inches(0.8), y + Inches(0.1), Inches(6.5), Inches(0.52),
        size=17, bold=True, color=RGBColor(0xfc, 0xa5, 0xa5))
    txt(s2, sub, Inches(0.8), y + Inches(0.65), Inches(6.5), Inches(0.5),
        size=13, color=MUTED)

for i, (title, sub) in enumerate(solutions):
    y = Inches(2.1 + i * 1.55)
    box(s2, Inches(8.5), y, Inches(7.0), Inches(1.3),
        color=RGBColor(0x04, 0x14, 0x0b), line=RGBColor(0x0e, 0x4a, 0x28))
    txt(s2, title, Inches(8.8), y + Inches(0.1), Inches(6.5), Inches(0.52),
        size=17, bold=True, color=GREEN)
    txt(s2, sub, Inches(8.8), y + Inches(0.65), Inches(6.5), Inches(0.5),
        size=13, color=MUTED)

txt(s2, "→", Inches(7.6), Inches(4.15), Inches(0.8), Inches(0.8),
    size=32, color=PURPLE, align=PP_ALIGN.CENTER, bold=True)

slide_num(s2, 2)


# ════════════════════════════════════════════════════════════
#  SLIDE 03 — 솔루션 개요
# ════════════════════════════════════════════════════════════
s3 = new_slide()
accent_bar(s3, GOLD)
section_tag(s3, "솔루션", GOLD)
heading(s3, "AI Dungeon RPG", GOLD)

cards = [
    ("◈", "완전 생성형 스토리", PURPLE,
     "플레이어의 모든 행동에\nAI가 고유하게 반응"),
    ("⚔", "풀 RPG 시스템", GOLD,
     "레벨업 · 인벤토리 · 퀘스트\nNPC · 하드코어 모드"),
    ("📖", "한국어 웹소설 문체", GREEN,
     "판타지 · SF · 공포 · 현대\n장르별 맞춤 문체"),
    ("♪", "장면 반응 BGM + SFX", BLUE,
     "전투·탐험·레벨업마다\nWeb Audio 자동 전환"),
]

for i, (icon, title, color, body) in enumerate(cards):
    l = Inches(0.5 + i * 3.85)
    bar = s3.shapes.add_shape(1, l, Inches(2.1), Inches(3.4), Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    box(s3, l, Inches(2.2), Inches(3.4), Inches(5.4), color=SURFACE, line=BORDER)
    txt(s3, icon, l, Inches(2.4), Inches(3.4), Inches(0.9),
        size=36, align=PP_ALIGN.CENTER)
    txt(s3, title, l, Inches(3.4), Inches(3.4), Inches(0.65),
        size=17, bold=True, color=color, align=PP_ALIGN.CENTER)
    multi_line(s3, body.split('\n'),
        l + Inches(0.2), Inches(4.2), Inches(3.0), Inches(2.0),
        size=14, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.6)

txt(s3, "▶  ai-dungeon-rpg.vercel.app",
    0, H - Inches(0.7), W, Inches(0.5),
    size=14, color=PURPLE, align=PP_ALIGN.CENTER, bold=True)

slide_num(s3, 3)


# ════════════════════════════════════════════════════════════
#  SLIDE 04 — 시스템 아키텍처
# ════════════════════════════════════════════════════════════
s4 = new_slide()
accent_bar(s4, PURPLE)
section_tag(s4, "아키텍처")
heading(s4, "시스템 구조")

nodes = [
    (Inches(0.4),  Inches(3.2), Inches(2.2), Inches(1.4),
     "👤 플레이어", "브라우저", PURPLE),
    (Inches(3.2),  Inches(2.4), Inches(3.2), Inches(3.0),
     "⚛ Frontend", "React + Vite\nZustand Store\nSSE Client\nVercel 배포", PURPLE),
    (Inches(6.8),  Inches(2.4), Inches(3.2), Inches(3.0),
     "⚡ Backend", "FastAPI\nSQLAlchemy\nJWT / OAuth\nRailway 배포", GREEN),
    (Inches(10.4), Inches(2.4), Inches(3.2), Inches(3.0),
     "🤖 Groq AI", "Llama 3.3 70B\nSSE 스트리밍\nFew-shot 프롬프트", GOLD),
]

for l, t, w, h, title, body, color in nodes:
    box(s4, l, t, w, h, color=RGBColor(0x0d, 0x0d, 0x1e), line=color)
    txt(s4, title, l, t + Inches(0.15), w, Inches(0.55),
        size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    multi_line(s4, body.split('\n'),
        l + Inches(0.15), t + Inches(0.75), w - Inches(0.3), h - Inches(0.9),
        size=12, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.4)

for x in [Inches(2.7), Inches(6.3), Inches(9.9)]:
    txt(s4, "→", x, Inches(3.55), Inches(0.7), Inches(0.6),
        size=20, color=BORDER, align=PP_ALIGN.CENTER, bold=True)

# DB
box(s4, Inches(6.8), Inches(6.0), Inches(3.2), Inches(0.95),
    color=SURFACE, line=BORDER)
txt(s4, "🗄 SQLite DB  (users · games · histories)",
    Inches(6.8), Inches(6.12), Inches(3.2), Inches(0.55),
    size=12, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
txt(s4, "↕", Inches(7.95), Inches(5.47), Inches(0.6), Inches(0.5),
    size=16, color=BORDER, align=PP_ALIGN.CENTER)

# SSE 라벨
txt(s4, "HTTP / SSE 스트리밍",
    Inches(3.2), Inches(1.85), Inches(7.3), Inches(0.45),
    size=11, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

slide_num(s4, 4)


# ════════════════════════════════════════════════════════════
#  SLIDE 05 — AI GM 핵심 기술
# ════════════════════════════════════════════════════════════
s5 = new_slide()
accent_bar(s5, PURPLE)
section_tag(s5, "핵심 기술")
heading(s5, "AI 게임마스터 설계")

features = [
    ("🔥 실시간 스트리밍",
     "SSE(Server-Sent Events)\nAI 토큰 생성 즉시 화면 출력\n타이프라이터 효과", PURPLE),
    ("📖 장르별 문체",
     "세계관 키워드로 자동 감지\n판타지·SF·공포·현대\n각기 다른 웹소설 서술체", GOLD),
    ("🎭 Few-shot 프리필링",
     "나쁜 예 + 좋은 예 쌍으로 주입\n언어 일관성 확보\n한자·영어 혼입 방지", GREEN),
    ("📦 구조화된 출력",
     "응답 마지막 JSON 블록 필수\nHP변화·아이템·퀘스트 추출\n서버가 자동 상태 갱신", BLUE),
]

for i, (title, body, color) in enumerate(features):
    col = i % 2
    row = i // 2
    l = Inches(0.5 + col * 7.85)
    t = Inches(2.1 + row * 3.0)
    bar = s5.shapes.add_shape(1, l, t, Inches(7.15), Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    box(s5, l, t + Inches(0.1), Inches(7.15), Inches(2.65), color=SURFACE, line=BORDER)
    txt(s5, title, l + Inches(0.25), t + Inches(0.25), Inches(6.7), Inches(0.6),
        size=18, bold=True, color=color)
    multi_line(s5, body.split('\n'),
        l + Inches(0.25), t + Inches(0.9), Inches(6.7), Inches(1.65),
        size=14, color=MUTED, spacing=1.5)

slide_num(s5, 5)


# ════════════════════════════════════════════════════════════
#  SLIDE 06 — 시연 영상 (중간, 핵심)
# ════════════════════════════════════════════════════════════
s6 = new_slide()
accent_bar(s6, GOLD)
section_tag(s6, "LIVE DEMO", GOLD)

txt(s6, "시연 영상",
    0, Inches(0.55), W, Inches(0.75),
    size=28, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

# 메인 플레이어 박스
VL, VT, VW, VH = Inches(1.2), Inches(1.6), Inches(13.6), Inches(5.9)
box(s6, VL, VT, VW, VH,
    color=RGBColor(0x06, 0x06, 0x12), line=GOLD, line_w=2)

txt(s6, "▶", VL + Inches(6.2), VT + Inches(1.2), Inches(1.2), Inches(1.2),
    size=64, color=GOLD, align=PP_ALIGN.CENTER)

txt(s6, "[ 시연 영상 ]",
    VL, VT + Inches(2.65), VW, Inches(0.6),
    size=15, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

multi_line(s6, [
    "로그인 → 세계관 선택 → 클래스 선택 → AI 오프닝 생성",
    "행동 입력 → GM 스트리밍 응답 → HP 플래시 + BGM 전환",
    "레벨업 → SFX + 스탯 애니메이션 → 캐릭터 시트 팝업",
], VL, VT + Inches(3.4), VW, Inches(1.9),
    size=14, color=MUTED, align=PP_ALIGN.CENTER, spacing=1.7)

slide_num(s6, 6)


# ════════════════════════════════════════════════════════════
#  SLIDE 07 — 게임 시스템
# ════════════════════════════════════════════════════════════
s7 = new_slide()
accent_bar(s7, GOLD)
section_tag(s7, "게임 시스템", GOLD)
heading(s7, "완전한 RPG 시스템")

systems = [
    ("⚔  캐릭터 & 클래스", "20종 직업별 고유 스탯\nHP·MP·힘·지능·민첩·카리스마"),
    ("📈  성장 시스템",     "전투 XP → 레벨업 → 스탯 보너스\nXP 가이드: 탐색 0~5 / 강적 80~150"),
    ("🎒  인벤토리",        "아이템 획득·소모·손실\n사망 패널티: 랜덤 아이템 손실"),
    ("📜  퀘스트 & NPC",    "동적 퀘스트 추가·완료 추적\nNPC 태도·설명 누적 기억"),
    ("🗺  방문 장소",       "visited_locations 자동 누적\n지도 탭에서 확인 가능"),
    ("💀  하드코어 모드",   "HP=0 → 서버가 강제 종료\n캐릭터 영구 사망 보장"),
]

for i, (title, body) in enumerate(systems):
    col = i % 3
    row = i // 3
    l = Inches(0.5 + col * 5.15)
    t = Inches(2.1 + row * 2.95)
    box(s7, l, t, Inches(4.7), Inches(2.6), color=SURFACE, line=BORDER)
    txt(s7, title, l + Inches(0.2), t + Inches(0.15), Inches(4.3), Inches(0.6),
        size=16, bold=True, color=GOLD)
    multi_line(s7, body.split('\n'),
        l + Inches(0.2), t + Inches(0.85), Inches(4.3), Inches(1.55),
        size=13, color=MUTED, spacing=1.45)

slide_num(s7, 7)


# ════════════════════════════════════════════════════════════
#  SLIDE 08 — Phase 2 UX·AI 개선
# ════════════════════════════════════════════════════════════
s8 = new_slide()
accent_bar(s8, PURPLE)
section_tag(s8, "Phase 2 개선")
heading(s8, "UX · AI 품질 대폭 향상")

ux_items = [
    "✦  타이프라이터 효과 (requestAnimationFrame)",
    "✦  스탯 플래시 애니메이션 (HP↑ 초록 / HP↓ 빨강)",
    "✦  StatusPanel 4탭 (상태 / 퀘스트 / NPC / 지도)",
    "✦  GM 텍스트 마크다운 하이라이트 (Bold → 골드)",
    "✦  SFX 효과음 4종 (전투 / 레벨업 / 아이템 / 오류)",
    "✦  키보드 단축키 (Esc 모달 닫기)",
    "✦  스토리 내보내기 (.txt 다운로드)",
    "✦  Dashboard 정렬·필터 (최신/턴/레벨 · 상태별)",
    "✦  캐릭터 생성 배경 스토리 프리셋",
    "✦  모험 갤러리 /stories (장르 필터)",
]

ai_items = [
    "◈  인벤토리·퀘스트 섹션 GM 프롬프트 삽입",
    "◈  NPC 일관성 규칙 강화",
    "◈  XP 가이드 수치화 (탐색 0~5 / 강적 80~150)",
    "◈  공개 스토리 API  GET /games/public",
    "◈  HttpError 클래스 — 4xx 재시도 방지",
    "◈  SQLite pool_pre_ping 안정성",
    "◈  token_count None 방어 처리",
]

box(s8, Inches(0.5), Inches(2.05), Inches(7.7), Inches(6.3),
    color=RGBColor(0x0a, 0x0a, 0x1c), line=PURPLE)
txt(s8, "FRONTEND  UX", Inches(0.75), Inches(2.15), Inches(7.2), Inches(0.5),
    size=12, bold=True, color=PURPLE)
multi_line(s8, ux_items,
    Inches(0.75), Inches(2.75), Inches(7.2), Inches(5.3),
    size=13, color=MUTED, spacing=1.55)

box(s8, Inches(8.4), Inches(2.05), Inches(7.1), Inches(6.3),
    color=RGBColor(0x0a, 0x12, 0x0a), line=GREEN)
txt(s8, "BACKEND  /  AI", Inches(8.65), Inches(2.15), Inches(6.6), Inches(0.5),
    size=12, bold=True, color=GREEN)
multi_line(s8, ai_items,
    Inches(8.65), Inches(2.75), Inches(6.6), Inches(5.3),
    size=13, color=MUTED, spacing=1.55)

slide_num(s8, 8)


# ════════════════════════════════════════════════════════════
#  SLIDE 09 — 챗봇 비교
# ════════════════════════════════════════════════════════════
s9 = new_slide()
accent_bar(s9, PURPLE)
section_tag(s9, "비교 분석")
heading(s9, "기존 AI 챗봇과의 차이")

rows = [
    ("캐릭터 상태 추적",  "매번 직접 설명",          "서버 자동 저장 & 갱신"),
    ("HP / XP 계산",      "사용자가 직접",            "서버 자동 처리"),
    ("장기 플레이 맥락",  "컨텍스트 초과 시 망각",    "자동 요약 압축"),
    ("게임 저장",         "없음 (대화 기록뿐)",       "DB 영구 저장"),
    ("하드코어 모드",     "규칙 합의 필요",           "시스템 수준 보장"),
    ("한국 웹소설 문체",  "일반 문어체",              "전용 프롬프트 + Few-shot"),
    ("BGM / SFX",         "없음",                     "Web Audio API 절차적 생성"),
    ("모험 공유",         "없음",                     "공개 링크 + 갤러리"),
]

# 헤더
hY = Inches(2.0)
box(s9, Inches(0.5), hY, Inches(5.2), Inches(0.55),
    color=RGBColor(0x14, 0x0a, 0x28), line=PURPLE)
txt(s9, "항목", Inches(0.7), hY + Inches(0.08), Inches(4.8), Inches(0.42),
    size=13, bold=True, color=PURPLE)
box(s9, Inches(5.8), hY, Inches(4.5), Inches(0.55),
    color=RGBColor(0x14, 0x08, 0x08), line=RED)
txt(s9, "ChatGPT / Claude", Inches(6.0), hY + Inches(0.08), Inches(4.1), Inches(0.42),
    size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
box(s9, Inches(10.4), hY, Inches(5.1), Inches(0.55),
    color=RGBColor(0x04, 0x14, 0x0b), line=GREEN)
txt(s9, "AI Dungeon RPG", Inches(10.6), hY + Inches(0.08), Inches(4.7), Inches(0.42),
    size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

for i, (item, bad, good) in enumerate(rows):
    rY = Inches(2.65 + i * 0.77)
    bg = RGBColor(0x0d, 0x0d, 0x1a) if i % 2 == 0 else SURFACE
    box(s9, Inches(0.5),  rY, Inches(5.2), Inches(0.7), color=bg, line=BORDER)
    box(s9, Inches(5.8),  rY, Inches(4.5), Inches(0.7), color=bg, line=BORDER)
    box(s9, Inches(10.4), rY, Inches(5.1), Inches(0.7), color=bg, line=BORDER)
    txt(s9, item, Inches(0.7),  rY + Inches(0.12), Inches(4.8), Inches(0.5),
        size=13, color=TEXT)
    txt(s9, bad,  Inches(6.0),  rY + Inches(0.12), Inches(4.1), Inches(0.5),
        size=12, color=RGBColor(0xfc, 0xa5, 0xa5), align=PP_ALIGN.CENTER)
    txt(s9, good, Inches(10.6), rY + Inches(0.12), Inches(4.7), Inches(0.5),
        size=12, color=GREEN, align=PP_ALIGN.CENTER)

slide_num(s9, 9)


# ════════════════════════════════════════════════════════════
#  SLIDE 10 — 기술적 도전과 해결
# ════════════════════════════════════════════════════════════
s10 = new_slide()
accent_bar(s10, RED)
section_tag(s10, "기술적 도전", RED)
heading(s10, "문제와 해결")

challenges = [
    ("한자 · 영어 혼입",
     "Llama 모델이 한국어 도중 微笑, こちらを 삽입",
     "장르별 프롬프트에 틀린 예/좋은 예 나란히 제시\ntext_sanitizer 후처리로 비한글 문자 제거"),
    ("JSON 스트리밍 노출",
     "상태 변화 JSON이 실시간 텍스트에 그대로 출력",
     "```json 감지 시점부터 텍스트 즉시 차단\nstripJson()로 최종 정리"),
    ("컨텍스트 폭증",
     "대화 누적 → 토큰 한도 초과 → GM 맥락 소실",
     "context_manager: 8,000토큰 초과 시 자동 요약\n최근 10턴 + 요약만 모델에 전달"),
]

for i, (title, prob, sol) in enumerate(challenges):
    t = Inches(2.1 + i * 2.2)
    box(s10, Inches(0.5), t, Inches(7.0), Inches(1.95),
        color=RGBColor(0x16, 0x07, 0x07), line=RGBColor(0x60, 0x18, 0x18))
    txt(s10, f"⚠  {title}", Inches(0.75), t + Inches(0.1), Inches(6.5), Inches(0.55),
        size=17, bold=True, color=RED)
    txt(s10, prob, Inches(0.75), t + Inches(0.72), Inches(6.5), Inches(1.0),
        size=13, color=MUTED)

    txt(s10, "→", Inches(7.65), t + Inches(0.65), Inches(0.9), Inches(0.65),
        size=26, color=PURPLE, align=PP_ALIGN.CENTER, bold=True)

    box(s10, Inches(8.7), t, Inches(6.8), Inches(1.95),
        color=RGBColor(0x04, 0x12, 0x08), line=RGBColor(0x0e, 0x4a, 0x28))
    txt(s10, "✅  해결", Inches(8.95), t + Inches(0.1), Inches(6.3), Inches(0.55),
        size=17, bold=True, color=GREEN)
    multi_line(s10, sol.split('\n'),
        Inches(8.95), t + Inches(0.72), Inches(6.3), Inches(1.1),
        size=13, color=MUTED, spacing=1.4)

slide_num(s10, 10)


# ════════════════════════════════════════════════════════════
#  SLIDE 11 — 수업 내용 연계
# ════════════════════════════════════════════════════════════
s11 = new_slide()
accent_bar(s11, GREEN)
section_tag(s11, "수업 연계", GREEN)
heading(s11, "수업에서 배운 내용의 적용")

connections = [
    ("LLM · 프롬프트 엔지니어링",
     GREEN,
     "모델이 토큰 단위로 생성한다는 원리 이해\nFew-shot: 나쁜 예/좋은 예 쌍으로 행동 유도\n구조화된 출력 형식 제어 (JSON 블록 필수화)"),
    ("API · 비동기 처리",
     PURPLE,
     "FastAPI (5주차)로 백엔드 전체 구축\nasync/await 비동기 SSE 스트리밍 구현\nSSE: 단방향 실시간 서버 푸시"),
    ("데이터 구조 · 상태 관리",
     GOLD,
     "캐릭터 상태 JSON 직렬화 → DB 저장\n슬라이딩 윈도우 알고리즘 → 컨텍스트 관리\nZustand: 프론트엔드 전역 상태 트리"),
    ("NLP · 텍스트 생성",
     BLUE,
     "프롬프트로 스타일 전이 (Style Transfer) 구현\n장르별 문체 제어 → 한국 웹소설 느낌\ntext_sanitizer: 정규식 후처리 파이프라인"),
]

for i, (title, color, body) in enumerate(connections):
    col = i % 2
    row = i // 2
    l = Inches(0.5 + col * 7.85)
    t = Inches(2.1 + row * 3.1)
    bar2 = s11.shapes.add_shape(1, l, t, Inches(7.15), Inches(0.08))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = color
    bar2.line.fill.background()
    box(s11, l, t + Inches(0.08), Inches(7.15), Inches(2.8), color=SURFACE, line=BORDER)
    txt(s11, title, l + Inches(0.25), t + Inches(0.2), Inches(6.7), Inches(0.6),
        size=17, bold=True, color=color)
    multi_line(s11, body.split('\n'),
        l + Inches(0.25), t + Inches(0.88), Inches(6.7), Inches(1.85),
        size=13, color=MUTED, spacing=1.5)

slide_num(s11, 11)


# ════════════════════════════════════════════════════════════
#  SLIDE 12 — 성과
# ════════════════════════════════════════════════════════════
s12 = new_slide()
accent_bar(s12, GOLD)
section_tag(s12, "성과", GOLD)
heading(s12, "구현 성과")

stats = [
    ("46",  "백엔드 테스트 통과",  PURPLE),
    ("20+", "직업 클래스",          GOLD),
    ("4",   "장르 지원",            GREEN),
    ("13+", "API 엔드포인트",       BLUE),
]

for i, (num, label, color) in enumerate(stats):
    l = Inches(0.5 + i * 3.8)
    bar3 = s12.shapes.add_shape(1, l, Inches(2.0), Inches(3.3), Inches(0.1))
    bar3.fill.solid(); bar3.fill.fore_color.rgb = color
    bar3.line.fill.background()
    box(s12, l, Inches(2.1), Inches(3.3), Inches(2.4), color=SURFACE, line=BORDER)
    txt(s12, num, l, Inches(2.15), Inches(3.3), Inches(1.3),
        size=60, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s12, label, l, Inches(3.55), Inches(3.3), Inches(0.55),
        size=14, color=MUTED, align=PP_ALIGN.CENTER)

divider(s12, Inches(5.0), BORDER)

features = [
    "✦ Google OAuth 로그인",
    "✦ Polar.sh 결제 & 플랜 관리",
    "✦ SSE 실시간 스트리밍",
    "✦ 타이프라이터 StreamText",
    "✦ Web Audio BGM + SFX 4종",
    "✦ StatusPanel 4탭",
    "✦ 컨텍스트 자동 압축",
    "✦ NPC·세계관 누적 기억",
    "✦ 하드코어 모드",
    "✦ 모험 갤러리 & 공유 링크",
    "✦ Dashboard 정렬·필터",
    "✦ 모바일 반응형 UI",
]

col1, col2, col3 = features[:4], features[4:8], features[8:]
multi_line(s12, col1,
    Inches(0.5), Inches(5.25), Inches(4.9), Inches(2.8),
    size=13, color=MUTED, spacing=1.55)
multi_line(s12, col2,
    Inches(5.6), Inches(5.25), Inches(4.9), Inches(2.8),
    size=13, color=MUTED, spacing=1.55)
multi_line(s12, col3,
    Inches(10.7), Inches(5.25), Inches(4.9), Inches(2.8),
    size=13, color=MUTED, spacing=1.55)

slide_num(s12, 12)


# ════════════════════════════════════════════════════════════
#  SLIDE 13 — 마무리
# ════════════════════════════════════════════════════════════
s13 = new_slide()
accent_bar(s13, PURPLE)

glow2 = s13.shapes.add_shape(9, Inches(3.5), Inches(1.2), Inches(9), Inches(6.5))
glow2.fill.solid()
glow2.fill.fore_color.rgb = RGBColor(0x12, 0x08, 0x26)
glow2.line.fill.background()

txt(s13, "⚔  AI DUNGEON RPG",
    0, Inches(1.6), W, Inches(0.65),
    size=13, color=PURPLE, align=PP_ALIGN.CENTER, bold=True)

txt(s13, "직접 플레이해보세요",
    0, Inches(2.4), W, Inches(1.1),
    size=50, bold=True, color=TEXT, align=PP_ALIGN.CENTER)

txt(s13, "단순한 데모가 아닌, 지금 이 순간에도 접속 가능한 라이브 서비스입니다",
    0, Inches(3.7), W, Inches(0.6),
    size=16, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

for i, (label, url, color) in enumerate([
    ("▶  플레이하기",  "ai-dungeon-rpg.vercel.app",               GOLD),
    ("⭐  GitHub",    "github.com/iringsoya-gif/ai-dungeon-rpg",   PURPLE),
]):
    l = Inches(3.2 + i * 5.0)
    box(s13, l, Inches(4.7), Inches(4.4), Inches(1.1),
        color=SURFACE, line=color, line_w=2)
    txt(s13, label, l, Inches(4.75), Inches(4.4), Inches(0.55),
        size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s13, url,   l, Inches(5.33), Inches(4.4), Inches(0.42),
        size=12, color=MUTED, align=PP_ALIGN.CENTER)

txt(s13, "감사합니다",
    0, Inches(6.5), W, Inches(0.85),
    size=30, bold=True, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

txt(s13, "202212577 박성안",
    0, H - Inches(0.6), W, Inches(0.45),
    size=12, color=MUTED, align=PP_ALIGN.CENTER)

slide_num(s13, 13)


# ── 저장 ──────────────────────────────────────────────────
OUTPUT = "presentation/AI_Dungeon_RPG_Presentation.pptx"
prs.save(OUTPUT)
print(f"✅  저장 완료: {OUTPUT}  ({TOTAL}슬라이드)")
